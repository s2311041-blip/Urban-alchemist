from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path(__file__).with_name("urban_alchemist_progress_2026-05.pptx")

FONT = "Hiragino Sans"
BG = RGBColor(247, 251, 255)
GREEN = RGBColor(15, 81, 50)
GREEN_2 = RGBColor(23, 98, 74)
DARK = RGBColor(23, 32, 51)
MUTED = RGBColor(78, 91, 108)
PALE_GREEN = RGBColor(226, 243, 233)
LIGHT_PANEL = RGBColor(255, 255, 255)
BORDER = RGBColor(185, 220, 205)
LIGHT_BORDER = RGBColor(215, 230, 222)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, left, top, width, height, text="", size=24, color=DARK, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    if align is not None:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle):
    add_textbox(slide, Inches(0.72), Inches(0.45), Inches(11.9), Inches(0.45), title, 23, GREEN_2, True)
    add_textbox(slide, Inches(0.72), Inches(0.92), Inches(12.0), Inches(0.8), subtitle, 31, DARK, True)
    line = slide.shapes.add_shape(SHAPE.RECTANGLE, Inches(0.72), Inches(1.82), Inches(2.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.color.rgb = GREEN


def add_flow(slide, labels, top=2.35):
    count = len(labels)
    left = 0.9 if count >= 4 else 1.6
    card_w = 2.55 if count >= 4 else 3.1
    card_h = 0.9
    for i, label in enumerate(labels):
        step = 3.05 if count >= 4 else 3.6
        x = Inches(left + i * step)
        card = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, x, Inches(top), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_PANEL
        card.line.color.rgb = BORDER
        tf = card.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = FONT
        r.font.bold = True
        r.font.size = Pt(17)
        r.font.color.rgb = GREEN

        if i < len(labels) - 1:
            offset = 2.64 if count >= 4 else 3.2
            add_textbox(slide, Inches(left + i * step + offset), Inches(top + 0.27), Inches(0.35), Inches(0.3), "→", 22, GREEN, True, PP_ALIGN.CENTER)


def add_icon_row(slide, icons, top=5.55):
    start_x = 1.45
    for i, (icon, label) in enumerate(icons):
        x = start_x + i * 3.65
        bubble = slide.shapes.add_shape(SHAPE.OVAL, Inches(x), Inches(top), Inches(0.75), Inches(0.75))
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = PALE_GREEN
        bubble.line.color.rgb = BORDER
        add_textbox(slide, Inches(x), Inches(top + 0.1), Inches(0.75), Inches(0.55), icon, 22, GREEN, True, PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.85), Inches(top + 0.18), Inches(2.45), Inches(0.4), label, 15, DARK, False)


def add_compare_panels(slide, left_title, left_points, right_title, right_points):
    panels = [
        (0.85, left_title, left_points),
        (6.85, right_title, right_points),
    ]
    for x, title, points in panels:
        panel = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.3), Inches(5.45), Inches(3.6))
        panel.fill.solid()
        panel.fill.fore_color.rgb = LIGHT_PANEL
        panel.line.color.rgb = LIGHT_BORDER
        add_textbox(slide, Inches(x + 0.25), Inches(2.52), Inches(5.0), Inches(0.5), title, 21, GREEN, True)

        content = "\n".join([f"• {line}" for line in points])
        tb = slide.shapes.add_textbox(Inches(x + 0.28), Inches(3.02), Inches(5.05), Inches(2.55))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.clear()
        lines = content.split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = FONT
            p.font.size = Pt(17)
            p.font.color.rgb = DARK
            p.space_after = Pt(8)


def add_three_cards(slide, cards):
    for i, (title, body) in enumerate(cards):
        x = 0.85 + i * 4.15
        panel = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.42), Inches(3.85), Inches(2.35))
        panel.fill.solid()
        panel.fill.fore_color.rgb = LIGHT_PANEL
        panel.line.color.rgb = LIGHT_BORDER
        add_textbox(slide, Inches(x + 0.22), Inches(2.63), Inches(3.4), Inches(0.45), title, 19, GREEN, True)
        add_textbox(slide, Inches(x + 0.22), Inches(3.08), Inches(3.4), Inches(1.5), body, 16, DARK, False)


def add_bullets(slide, items, top=5.1):
    tb = slide.shapes.add_textbox(Inches(0.95), Inches(top), Inches(11.8), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.font.name = FONT
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_note(slide, note):
    shape = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.2), Inches(11.4), Inches(0.58))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE_GREEN
    shape.line.color.rgb = PALE_GREEN
    frame = shape.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = note
    run.font.name = FONT
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = GREEN


def add_footer(slide, idx):
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(0.7), Inches(0.25), str(idx), 11, MUTED, False, PP_ALIGN.RIGHT)


def slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, Inches(0.88), Inches(1.55), Inches(11.6), Inches(0.75), "ゼミ進捗報告", 34, GREEN_2, True)
    add_textbox(slide, Inches(0.88), Inches(2.45), Inches(11.6), Inches(0.75), "Urban Alchemist プロトタイプ開発", 38, DARK, True)
    name_box = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.8), Inches(6.6), Inches(0.8))
    name_box.fill.solid()
    name_box.fill.fore_color.rgb = PALE_GREEN
    name_box.line.color.rgb = PALE_GREEN
    add_textbox(slide, Inches(1.15), Inches(4.02), Inches(6.0), Inches(0.4), "発表者: あなたの名前（ここを置き換え）", 22, GREEN, True)
    add_textbox(slide, Inches(0.9), Inches(4.95), Inches(6.6), Inches(0.3), "所属・日付（必要なら追記）", 15, MUTED, False)
    add_footer(slide, 1)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, Inches(0.82), Inches(1.2), Inches(11.7), Inches(0.9), "Urban Alchemist", 47, GREEN, True)
    add_textbox(slide, Inches(0.85), Inches(2.05), Inches(11.6), Inches(0.9), "都市課題を「集める」から「改善まで回す」へ", 31, DARK, True)
    tag = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(3.35), Inches(7.3), Inches(0.58))
    tag.fill.solid()
    tag.fill.fore_color.rgb = PALE_GREEN
    tag.line.color.rgb = PALE_GREEN
    add_textbox(slide, Inches(1.05), Inches(3.5), Inches(7.0), Inches(0.3), "AR投稿 × VR/ゲーム体験 × 参加型まちづくり", 18, GREEN, True)
    add_flow(slide, ["課題を投稿", "ゲームで体験", "改善案を実装", "次の参加へ"], top=4.5)
    add_footer(slide, 2)


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "今日のアジェンダ", "今回の報告で伝えること")
    topics = [
        ("1. 背景と問題意識", ["なぜゲーム化が必要か", "前回までとの違い"]),
        ("2. 提案と独自性", ["投稿を体験可能な課題へ変換", "継続参加と多様な視点"]),
        ("3. 実装デモ", ["不満発見から改善までを実演", "島の変化を確認"]),
        ("4. 今後の計画", ["次回までの実施項目", "検証の進め方"]),
    ]
    for idx, (title, lines) in enumerate(topics):
        col = idx % 2
        row = idx // 2
        x = 0.9 + col * 6.1
        y = 2.2 + row * 2.0
        panel = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.6), Inches(1.75))
        panel.fill.solid()
        panel.fill.fore_color.rgb = LIGHT_PANEL
        panel.line.color.rgb = LIGHT_BORDER
        add_textbox(slide, Inches(x + 0.2), Inches(y + 0.15), Inches(5.2), Inches(0.35), title, 20, GREEN, True)
        add_textbox(slide, Inches(x + 0.25), Inches(y + 0.55), Inches(5.1), Inches(1.05), f"• {lines[0]}\n• {lines[1]}", 16, DARK, False)
    add_footer(slide, 3)


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "問題意識", "なぜ「ゲーム化」が必要か")
    add_flow(slide, ["投稿の手間", "単発参加で終わる", "継続データが不足"], top=2.45)
    add_bullets(
        slide,
        [
            "参加を一回で終わらせない仕組みが必要",
            "都市課題を「自分ごと」に変える体験が必要",
            "その役割を担うのがゲームのループ",
        ],
        top=4.6,
    )
    add_footer(slide, 4)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "提案する体験", "プレイヤーが「改善の当事者」になる")
    add_flow(slide, ["🔍 不満を発見", "🧠 改善方法を選ぶ", "🛠 建築して解決", "🌱 まちが変化"], top=2.32)
    add_compare_panels(
        slide,
        "プレイヤーが行うこと",
        ["不満の内容を読む", "改善プランを選択する", "建築して効果を確認する"],
        "研究として得られること",
        ["どの課題が選ばれたか", "どの改善策が選ばれたか", "改善後に行動が続くか"],
    )
    add_icon_row(
        slide,
        [("🧍", "探索する"), ("🧱", "建築する"), ("✅", "改善を確認")],
        top=5.52,
    )
    add_note(slide, "都市課題を「見る」だけでなく、直す行為と継続行動まで扱える")
    add_footer(slide, 5)


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "独自性 1", "投稿データを「体験可能な課題」に変える")
    add_compare_panels(
        slide,
        "一般的な流れ",
        ["投稿 → DB保存", "専門家が分析", "市民の関与が止まりやすい"],
        "本研究の流れ",
        ["投稿 → ゲーム課題化", "他者が体験", "改善案を作り参加が連鎖"],
    )
    add_footer(slide, 6)


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "独自性 2", "継続参加を設計で作る")
    add_three_cards(
        slide,
        [
            ("内発的動機", "自分のまちを良くする\n作る楽しさ"),
            ("外発的動機", "達成表示\n島の拡張"),
            ("循環設計", "解決すると次の探索へ\n自然に続く"),
        ],
    )
    add_note(slide, "「お願いして参加」ではなく「遊んで参加」へ")
    add_footer(slide, 7)


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "独自性 3", "多様な視点を取り入れる都市体験")
    add_bullets(
        slide,
        [
            "不満を「誰にとってのバリアか」で捉える",
            "暗さ・段差・案内不足・移動しづらさを体験化",
            "改善プランを比較し、誰に有効かを考える",
        ],
        top=2.4,
    )
    add_flow(slide, ["高齢者の視点", "子育ての視点", "車いすの視点"], top=4.55)
    add_note(slide, "バリア理解・UD・インクルーシブデザインを実感ベースで学べる")
    add_footer(slide, 8)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "今回の実装とデモ", "ここまで実装済み")
    add_bullets(
        slide,
        [
            "不満発見 → プラン選択 → 建築解決",
            "解決後の島拡張 / フェリー / 地図 / コンパス",
            "次ステップ: AR投稿接続、効果検証",
        ],
        top=2.5,
    )
    add_flow(slide, ["1. 不満を開く", "2. 改善して完成", "3. 変化を確認"], top=4.8)
    add_note(slide, "デモ: 改善前→改善→島の変化を2〜3分で確認")
    add_footer(slide, 9)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "次回までの課題", "次の進捗で実施すること")
    add_three_cards(
        slide,
        [
            ("VR体験のデザイン", "投稿課題をどう体験化するかの導線設計"),
            ("論文調査", "参加型設計・ゲーミフィケーション・UD関連の整理"),
            ("ゲームのさらなる実装", "AR投稿連携に向けた機能追加と検証準備"),
        ],
    )
    add_note(slide, "次回は「実装進捗 + 検証設計」をセットで報告")
    add_footer(slide, 10)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build()
